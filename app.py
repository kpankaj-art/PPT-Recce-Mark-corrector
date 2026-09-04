import io
import json
import warnings
import numpy as np
import cv2
import streamlit as st
from PIL import Image, ImageDraw
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from google import genai

warnings.filterwarnings("ignore")

st.set_page_config(page_title="PPT Image Line Cleaner", page_icon="🎨", layout="centered")

st.title("🎨 PPT Hand-drawn Line Erase & Border Tool")
st.write("Photo ke andar ki rough markings ko detect karke unhe erase karega aur clean rectangular border banaye ka.")

GEMINI_API_KEY = st.secrets.get("GEMINI_API_KEY", "")

@st.cache_resource
def get_client(api_key):
    return genai.Client(api_key=api_key)

if not GEMINI_API_KEY:
    st.error("Streamlit Secrets mein GEMINI_API_KEY missing hai!")
    st.stop()

client = get_client(GEMINI_API_KEY)

def detect_marking_area(pil_image):
    """Gemini Flash se image ke andar ki rough line/annotation ka bounding box aur color nikalna"""
    prompt = """
    Analyze this image for hand-drawn irregular lines, circles, scribble marks, or rough freehand boxes over shop boards/objects.
    Respond strictly in JSON format:
    {
      "has_mark": true,
      "color": [R, G, B],
      "box_percent": [ymin, xmin, ymax, xmax]
    }
    where box_percent contains percentage coordinates (0 to 100) tightly surrounding the hand-drawn mark.
    """
    try:
        response = client.models.generate_content(
            model='gemini-1.5-flash',
            contents=[pil_image, prompt]
        )
        clean_text = response.text.replace("```json", "").replace("```", "").strip()
        data = json.loads(clean_text)

        if not data.get("has_mark", False):
            return None, None

        box = data["box_percent"]
        color = tuple(data.get("color", [255, 0, 0]))
        return box, color
    except Exception:
        return None, None

def remove_mark_and_draw_rect(pil_image, box, color):
    """Computer Vision Inpainting: Rough merged drawing ko erase karke clean rectangle draw karna"""
    w, h = pil_image.size
    ymin = int((box[0] / 100) * h)
    xmin = int((box[1] / 100) * w)
    ymax = int((box[2] / 100) * h)
    xmax = int((box[3] / 100) * w)

    # Convert PIL Image to OpenCV Format (BGR)
    img_np = np.array(pil_image.convert("RGB"))
    img_cv = cv2.cvtColor(img_np, cv2.COLOR_RGB2BGR)

    # Create mask for the bounding box boundary (where hand drawn lines lie)
    mask = np.zeros((h, w), dtype=np.uint8)
    
    # 10px thick boundary line mask around detected object
    thickness = 12
    cv2.rectangle(
        mask, 
        (max(0, xmin - thickness), max(0, ymin - thickness)), 
        (min(w, xmax + thickness), min(h, ymax + thickness)), 
        255, 
        thickness=thickness*2
    )

    # OpenCV Inpainting: Mitao rough lines by blending background pixels
    inpainted_cv = cv2.inpaint(img_cv, mask, inpaintRadius=5, flags=cv2.INPAINT_TELEA)

    # Convert back to PIL Image
    inpainted_rgb = cv2.cvtColor(inpainted_cv, cv2.COLOR_BGR2RGB)
    cleaned_pil = Image.fromarray(inpainted_rgb)

    # Draw fresh straight rectangular border over cleaned area
    draw = ImageDraw.Draw(cleaned_pil)
    draw.rectangle([xmin, ymin, xmax, ymax], outline=color, width=4)

    return cleaned_pil

def process_presentation(uploaded_file):
    prs = Presentation(uploaded_file)
    total_slides = len(prs.slides)
    
    progress_bar = st.progress(0)
    status_text = st.empty()

    cleaned_count = 0

    for idx, slide in enumerate(prs.slides, start=1):
        status_text.text(f"Processing Slide {idx}/{total_slides}...")
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    img_bytes = shape.image.blob
                    img_pil = Image.open(io.BytesIO(img_bytes)).convert("RGB")

                    # 1. Detection
                    box, color = detect_marking_area(img_pil)
                    
                    if box:
                        # 2. Erase rough mark & redraw clean box
                        cleaned_pil = remove_mark_and_draw_rect(img_pil, box, color)

                        # 3. Save back into PPTX shape
                        output = io.BytesIO()
                        cleaned_pil.save(output, format="PNG")
                        shape.image.blob = output.getvalue()
                        cleaned_count += 1
                except Exception:
                    pass
        progress_bar.progress(idx / total_slides)

    status_text.text(f"Processing Complete! Cleaned {cleaned_count} images.")
    
    ppt_out = io.BytesIO()
    prs.save(ppt_out)
    ppt_out.seek(0)
    return ppt_out, cleaned_count

uploaded_file = st.file_uploader("Upload PowerPoint File (.pptx)", type=["pptx"])

if uploaded_file is not None:
    if st.button("Fix & Clean PPT", type="primary"):
        with st.spinner("AI photo me se rough lines erase kar ke straight box draw kar raha hai..."):
            fixed_ppt, count = process_presentation(uploaded_file)
            
            if count > 0:
                st.success(f"Success! Total {count} images fix ho gayi hain.")
            else:
                st.info("Koi hand-drawn mark detect nahi hua ya PPT image format different hai.")

            st.download_button(
                label="📥 Download Fixed PPT",
                data=fixed_ppt,
                file_name=f"fixed_{uploaded_file.name}",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
            )
